"""Create a PowerPoint report"""

# Standard library imports
import functools
import io

# Third party imports
import pptx
import pyplugs
from pptx.chart.data import ChartData
from pptx.enum import chart as pptx_charts
from pyconfs import Configuration

# Geo:N:G imports
from geong_common import files

# Call plugin functions within powerpoint module
*_, PLUGIN = __name__.split(".")
call = functools.partial(pyplugs.call, __package__, PLUGIN)
content_types = functools.partial(pyplugs.funcs, __package__, PLUGIN)


@pyplugs.register
def generate_powerpoint(template_url, cfg, data, output_path):
    """Generate PowerPoint report from data based on cfg. Store to output_path"""
    template = files.get_url_or_asset(template_url, local_assets="geong_common.assets")
    template_cfg = Configuration.from_str(template.read_text(), format="toml")
    template_ppt = (template.parent / template_cfg.replace("template")).read_bytes()

    prs = pptx.Presentation(io.BytesIO(template_ppt))
    add_slides(prs, template_cfg, cfg, data)
    prs.save(output_path)


def add_slides(prs, template, cfg, data):
    """Generate slides and insert them into the prs PowerPoint document"""
    for slide_cfg in cfg.slide:
        if "nested" in slide_cfg:
            for single_data in data[slide_cfg.nested["data"]]:
                add_slides(prs, template, cfg[slide_cfg.nested.cfg], single_data)
            continue

        layout_name = template.layouts.get(slide_cfg.layout, slide_cfg.layout)
        layout = prs.slide_layouts.get_by_name(layout_name)
        slide = prs.slides.add_slide(layout)

        # Add title
        if "title" in slide_cfg:
            slide.shapes.title.text = slide_cfg.title.format(**data)

        # Add content
        ph_cfg = template.placeholders[slide_cfg.layout]
        for placeholder, info in slide_cfg.section_items:
            if placeholder not in ph_cfg:
                raise ValueError(
                    f"Unknown placeholder {placeholder!r}. "
                    f"Choose between {', '.join(ph_cfg.entry_keys)}"
                )

            func = f"add_{info.type}"
            if func not in content_types():
                types = [ct[4:] for ct in content_types() if ct.startswith("add_")]
                raise ValueError(
                    f"Unknown type {info.type!r}. Choose between {', '.join(types)}"
                )

            content = info.content
            if isinstance(content, str) and content.startswith("ref:"):
                content = cfg.content[info.type][info.content[4:]]

            call(
                func=func,
                placeholder=slide.placeholders[ph_cfg[placeholder]],
                content=content,
                data=data,
            )


@pyplugs.register
def add_text(placeholder, content, data):
    """Add text to placeholder"""
    placeholder.text = content.format(**data)


@pyplugs.register
def add_picture(placeholder, content, data):
    """Add picture to placeholder, content should be path to image file"""
    placeholder.insert_picture(content.format(**data))


@pyplugs.register
def add_chart(placeholder, content, data):
    """Add chart to placeholder"""
    # Create a common index
    categories = []
    for col in content["data"]:
        for key in data[col]:
            if key not in categories:
                categories.append(key)

    # Add chart data
    chart = ChartData()
    chart.categories = categories
    for series_name, col in zip(content["series_names"], content["data"]):
        values = [data[col][cat] for cat in categories]
        chart.add_series(series_name, values)

    # Insert chart
    placeholder.insert_chart(
        getattr(pptx_charts.XL_CHART_TYPE, content.type.upper()), chart
    )


@pyplugs.register
def add_table(placeholder, content, data):
    """Add table to placeholder"""
    if "column_names" in content:
        _add_horisontal_table(placeholder, content, data)
    elif "row_names" in content:
        _add_vertical_table(placeholder, content, data)
    else:
        raise ValueError("Either 'column_names' or 'row_names' should be defined")


def _add_horisontal_table(placeholder, content, data):
    """Add horisontal table to placeholder"""
    # Create a common index
    index_names = []
    for col in content["data"]:
        for key in data[col]:
            if key not in index_names:
                index_names.append(key)

    # Set up column names
    column_names = [content["index_name"]] + content["column_names"]

    # Add table with headers
    table = placeholder.insert_table(
        rows=len(index_names) + 1, cols=len(column_names)
    ).table
    for idx, column_name in enumerate(column_names):
        table.cell(0, idx).text = column_name
    for idx, index_name in enumerate(index_names, start=1):
        table.cell(idx, 0).text = index_name

    # Add data to table
    for col_idx, (col, fmt) in enumerate(
        zip(content["data"], content["formats"]), start=1
    ):
        for row_idx, row in enumerate(index_names, start=1):
            if row not in data[col]:
                continue
            value = data[col][row]
            table.cell(row_idx, col_idx).text = fmt.format(value) if fmt else str(value)


def _add_vertical_table(placeholder, content, data):
    """Add vertical table to placeholder"""
    raise NotImplementedError("Vertical tables are not yet implemented")
