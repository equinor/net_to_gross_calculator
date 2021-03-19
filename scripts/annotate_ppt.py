"""Annotate a PowerPoint template

Based on https://pbpython.com/creating-powerpoint.html
"""

# Standard library imports
import pathlib

# Third party imports
import typer
from pptx import Presentation

# Geo:N:G imports
from geong_common import log
from geong_common.log import logger

log.init()


def cli():
    """Use Typer to create a Command Line Interface"""
    typer.run(analyze_ppt)


def analyze_ppt(input: pathlib.Path, output: pathlib.Path):
    """Analyze the structure of a PowerPoint template.

    The output file contains marked up information that can be used to specify
    TOML configuration files.
    """
    logger.info(f"Read template {input.resolve()}")
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, layout in enumerate(prs.slide_layouts):
        logger.info(f"Analyzing layout {index}: {layout.name}")
        slide = prs.slides.add_slide(layout)
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = f"Title for Layout {index}"
        except AttributeError:
            logger.warning(f"No Title for Layout {index}")
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = f"Placeholder index:{phf.idx} type:{shape.name}"
                except AttributeError:
                    logger.warning(f"{phf.type} has no text attribute")
                logger.info(f"- Placeholder index {phf.idx}: {shape.name}")

    logger.info(f"Save template information in {output.resolve()}")
    prs.save(output)


if __name__ == "__main__":
    cli()
